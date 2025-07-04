import io
import pandas as pd
import docx
import pptx
from pdfminer.high_level import extract_text
from pptx import Presentation

class IngestionAgent:
    def __init__(self):
        pass

    def parse_documents(self, uploaded_files):
        combined_text = ""
        for file in uploaded_files:
            filename = file.name.lower()
            if filename.endswith('.pdf'):
                combined_text += self.parse_pdf(file)
            elif filename.endswith('.docx'):
                combined_text += self.parse_docx(file)
            elif filename.endswith('.pptx'):
                combined_text += self.parse_pptx(file)
            elif filename.endswith('.csv'):
                combined_text += self.parse_csv(file)
            elif filename.endswith('.txt') or filename.endswith('.md'):
                combined_text += self.parse_text(file)
            combined_text += "\n\n"
        return combined_text

    def parse_pdf(self, file):
        return extract_text(file)

    def parse_docx(self, file):
        doc = docx.Document(file)
        return "\n".join([para.text for para in doc.paragraphs])

    def parse_pptx(self, file):
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def parse_csv(self, file):
        df = pd.read_csv(file)
        return df.to_string(index=False)

    def parse_text(self, file):
        return file.read().decode("utf-8")